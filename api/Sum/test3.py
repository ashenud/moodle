from pptx import Presentation

prs = Presentation()

lyt = prs.slide_layouts[0]  # choosing a slide layout

for x in range(0, 3):

    if x == 2:
        slide = prs.slides.add_slide(lyt)  # adding a slide
        title = slide.shapes.title  # assigning a title
        subtitle = slide.placeholders[1]  # placeholder for subtitle
        subtitle.text = "'Miss Brill' is the story of an old woman told brilliantly and realistically, balancing thoughts and emotions that sustain her late solitary life amidst all the bustle of modern life. Miss Brill is a regular visitor on Sundays to the Jardins Publiques (the Public Gardens) of a small French suburb where she sits and watches all sorts of people come and go. She listens to the band playing, loves to watch people and guess what keeps them going, and enjoys contemplating the world as a great stage upon which actors perform. She finds herself to be another actor among the so many she sees, or at least herself as 'part of the performance after all.' One Sunday Miss Brill puts on her fur and goes to the Public Gardens as usual. The evening ends with her sudden realization that she is old and lonely, a realization brought to her by a conversation she overhears between a boy and a girl, presumably lovers, who comment on her unwelcome presence in their vicinity. Miss Brill is sad and depressed as she returns home, not stopping by as usual to buy her Sunday delicacy, a slice of honey-cake. She retires to her dark room, puts the fur back into the box and imagines that she has heard something cry."  # subtitle
    else:
        slide = prs.slides.add_slide(lyt)  # adding a slide
        title = slide.shapes.title  # assigning a title
        subtitle = slide.placeholders[1]  # placeholder for subtitle
        title.text = "Hey,This is a Slide! How exciting!"  # title
        subtitle.text = "Really?"  # subtitle

prs.save("slide3.pptx")  # saving file
