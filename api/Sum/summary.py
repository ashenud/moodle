# importing libraries
import glob
from pptx import Presentation
import math
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
import textract
import os.path
#nltk.download()

def create_sumall(abc, ratio):
    # filename = "*.pptx"
    if abc:
        filename = abc
        stop_word = ['is', 'a', 'and', 'the']  # Words which needs to ignore from the slide

        # Function to create Text summarization
        def create_summ(text):
            stopWords = set(stopwords.words("english"))
            words = word_tokenize(text)
            freqTable = dict()
            for word in words:
                word = word.lower()
                if word in stopWords:
                    continue
                if word in freqTable:
                    freqTable[word] += 1
                else:
                    freqTable[word] = 1

            sentences = sent_tokenize(text)
            sentenceValue = dict()

            for sentence in sentences:
                for word, freq in freqTable.items():
                    if word in sentence.lower():
                        if sentence in sentenceValue:
                            sentenceValue[sentence] += freq
                        else:
                            sentenceValue[sentence] = freq

            sumValues = 0
            for sentence in sentenceValue:
                sumValues += sentenceValue[sentence]

            lensenvalu = len(sentenceValue)
            if lensenvalu == 0:
                lensenvalu = 1
                average = int(sumValues / lensenvalu)
            else:
                average = int(sumValues / lensenvalu)

            summary = ''
            for sentence in sentences:
                if (sentence in sentenceValue) and (sentenceValue[sentence] > (
                        ratio * average)):  # Based on requirement set the range between 1.0 - 1.2
                    summary += " " + sentence

            return summary

        # Reading Full pptx and passing the text into summarization function
        def read_full_pptxe(filename):
            sentences = []
            b = []
            a = 0
            for eachfile in glob.glob(filename):
                prs = Presentation(eachfile)
                for slide in prs.slides:
                    a = a + 1
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            s = create_summ(shape.text.replace("\n", " "))
                            s = str(s)
                            if (len(s)) >= 1:
                                f = ["Slide " + str(a) + "-" + s]
                                # b.append(f)
                                # print(f)
                                sentences.append(f)
            return sentences

        def read_full_docx(filename):
            sentences = []
            text = textract.process(filename)
            temp = text.split(".")
            for t in temp:
                sentences.append(t.replace("\n", " "))
            return sentences

        # Print slide summarization text with number
        extension = os.path.splitext(filename)[1]

        if extension == 'docx':
            read_full_docx(filename)
        else:
            read_full_pptxe(filename)

        # Function to convert string into list
        def Convert(string):
            li = list(string.split(" "))
            return li

        # Function to convert String with new line into list
        def Convert2(string):
            li = list(string.split("\n"))
            return li

        # Function to read only Slide 3 i.e. Learning Outcomes
        def read_slide3(filename):
            a = 1
            for eachfile in glob.glob(filename):
                prs = Presentation(eachfile)
                for slide in prs.slides:
                    a = a + 1
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if a == 4 and shape.shape_id == 3:
                                s3 = str(shape.text)
                                return s3

        # Function to read pptx and Match the slide number with Learning Outcomes point
        def read_full_pptx(filename, sss):
            numberslide = []
            numberslide.append(sss)
            a = 0
            # print("\n")
            # print(sss)
            # print('---------------------------------------------')
            for eachfile in glob.glob(filename):
                prs = Presentation(eachfile)
                for slide in prs.slides:
                    a = a + 1
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if shape.shape_id != 2:
                                s = shape.text.replace("\n", " ")
                                s = str(s)
                                if (len(s)) >= 20 and a != 3 and a != 1 and a != 2:
                                    lo_1 = [a for a in new_l1 if a in s.lower()]
                                    f_lo_l = round((len(lo_1) / len_of_l1) * 100)
                                    if f_lo_l >= 50:
                                        f = "Slide " + str(a)
                                        # print(f)
                                        numberslide.append(f)

            return numberslide

        loooo = Convert2(read_slide3(filename))

        # Loop to read Learning outcome point in line and calling read_full_pptx function to matching with whole slide
        abc = []
        for i in loooo:
            l1 = Convert(i.lower())
            new_l1 = [w for w in l1 if w not in stop_word]
            len_of_l1 = len(new_l1)
            read_full_pptx(filename, i)
            # print(read_full_pptx(filename,i))
            abc.append(read_full_pptx(filename, i))
        # print('dshfdsufgsdhfsdgf',abc)

        return (read_full_pptxe(filename), abc)
    else:
        print('error')
