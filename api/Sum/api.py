from flask import Flask, request, url_for, redirect, render_template
from flask_cors import CORS
import werkzeug
import summary as summarizeed
import json
import textract
from pptx import Presentation
import os

app = Flask(__name__)
CORS(app)


@app.route('/summerize', methods=['GET', 'POST'])
def summerize():
    file = request.files['file']
    ratio = float(request.form['ratio'])
    filename = werkzeug.utils.secure_filename(file.filename)
    print("\nReceived image File name : " + file.filename)
    file.save('upload/' + filename)

    f, file_extension = os.path.splitext('upload/' + filename)
    print(file_extension)

    if file_extension == '.docx':
        text = textract.process('upload/' + filename)
        arr = str(text).replace("\\n", "")
        arr = arr.replace("\\t", "")
        arr = arr.replace("\\", "")

        prs = Presentation()

        lyt = prs.slide_layouts[0]  # choosing a slide layout

        for x in range(0, 3):

            if x == 2:
                slide = prs.slides.add_slide(lyt)  # adding a slide
                title = slide.shapes.title  # assigning a title
                subtitle = slide.placeholders[1]  # placeholder for subtitle
                subtitle.text = arr
            else:
                slide = prs.slides.add_slide(lyt)  # adding a slide
                title = slide.shapes.title  # assigning a title
                subtitle = slide.placeholders[1]  # placeholder for subtitle
                title.text = "ignore"  # title
                subtitle.text = "ignore"  # subtitle
        prs.save("upload/slide3.pptx")  # saving file
        print('file saved')
        res = summarizeed.create_sumall('upload/slide3.pptx', ratio)
    else:
        res = summarizeed.create_sumall('upload/' + filename, ratio)

    rr = []
    for r in res[0]:
        rr.append(r[0].replace('"', ''))

    return_str = '{ "result" : ['

    for i in range(len(rr)):
        if i == len(rr) - 1:
            return_str += '"' + rr[i] + '"'
        else:
            return_str += '"' + rr[i] + '"' + ','
    return_str += ']}'

    print(return_str)

    return json.loads(return_str)


if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5001, debug=True)
